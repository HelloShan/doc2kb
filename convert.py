"""
doc2kb — 文档转换引擎模块
=============================
转换引擎优先级：
  主力：Docling（支持 docx/pdf/xlsx/pptx，高质量表格还原+多栏识别）
  降级：原生 Python 解析器（python-docx / pypdf / openpyxl / python-pptx）

支持将 6 种核心格式转换为清洁后的 Markdown：
  - .docx  → Docling(主) / python-docx(降)
  - .pdf   → Docling(主) / pypdf(降)
  - .xlsx  → Docling(主) / openpyxl(降)
  - .pptx  → Docling(主) / python-pptx(降)
  - .md    → 复制+乱码校验
  - .txt   → 直接复制，编码自动回退

每个文件返回 (status, md_rel_path_or_None, error_msg_or_None, warning_or_None) 四元组。
"""

import os
import re
import traceback
from pathlib import Path
from typing import Optional, Tuple, List
from concurrent.futures import ThreadPoolExecutor, as_completed

from config import (
    SOURCE_DIR, OUTPUT_MD_DIR, SUPPORTED_EXTENSIONS,
    CONVERT_WORKERS, MAX_MD_FILE_SIZE_KB,
)
from validate import is_file_readable
from state import compute_sha256
from docx.opc.exceptions import PackageNotFoundError


# ============================================================
# 工具函数
# ============================================================

def _get_rel_path(source_path: Path) -> str:
    return source_path.relative_to(SOURCE_DIR).as_posix()


def _get_output_md_path(source_path: Path) -> Path:
    rel = source_path.relative_to(SOURCE_DIR)
    return OUTPUT_MD_DIR / rel.with_suffix(".md")


def _ensure_output_dir(output_path: Path):
    output_path.parent.mkdir(parents=True, exist_ok=True)


# ============================================================
# Docling 引擎（主力）
# ============================================================

_DOCLING_CONVERTER = None
_DOCLING_SUPPORTED = {".docx", ".pdf"}


def _get_docling_converter():
    global _DOCLING_CONVERTER
    if _DOCLING_CONVERTER is None:
        from docling.document_converter import DocumentConverter
        _DOCLING_CONVERTER = DocumentConverter()
    return _DOCLING_CONVERTER


def _is_docm_file(source_path: Path) -> bool:
    """
    检测伪装为 .docx 的宏文档 (.docm)。
    从 ZIP 包中的 [Content_Types].xml 检查 ContentType 标记。
    """
    if source_path.suffix.lower() != '.docx':
        return False
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(source_path)) as z:
            ct = z.read('[Content_Types].xml')
            root = ET.fromstring(ct)
            ns = root.tag.split('}')[0].strip('{') if '}' in root.tag else ''
            tag = f'{{{ns}}}Override' if ns else 'Override'
            for override in root.iter(tag):
                ct_type = override.get('ContentType', '')
                if 'macroenabled' in ct_type.lower():
                    return True
        return False
    except Exception:
        return False


# ============================================================
# 兼容性检测（供 doc_pipeline.py check 子命令使用）
# ============================================================

def detect_docm(file_path: Path) -> tuple[bool, str]:
    """检测 .docx 文件是否为宏文档 (.docm) 或旧版 .doc。返回 (是问题吗, 原因)。"""
    if file_path.suffix.lower() not in ('.docx',):
        return False, ''
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(file_path)) as z:
            ct = z.read('[Content_Types].xml')
            root = ET.fromstring(ct)
            ns = root.tag.split('}')[0].strip('{') if '}' in root.tag else ''
            tag = f'{{{ns}}}Override' if ns else 'Override'
            for override in root.iter(tag):
                ct_type = override.get('ContentType', '')
                if 'macroenabled' in ct_type.lower():
                    return True, '宏文档 (.docm)，需用 Word 另存为 .docx'
        return False, ''
    except (zipfile.BadZipFile, Exception):
        return True, '不是有效的 ZIP 包（可能是旧版 .doc 格式误标为 .docx）'


def detect_broken_pdf(file_path: Path) -> tuple[bool, str]:
    """检查 PDF 文件是否可读（静默模式，压制 pypdf 的 stderr 和日志噪音）。"""
    import logging
    import contextlib

    if file_path.suffix.lower() != '.pdf':
        return False, ''
    try:
        from pypdf import PdfReader, errors as pypdf_errors
        # 静默 pypdf 的日志和 stderr 噪音（如 incorrect startxref pointer）
        logger = logging.getLogger('pypdf')
        old_level = logger.level
        logger.setLevel(logging.ERROR)
        try:
            devnull = open(os.devnull, 'w')
            try:
                with contextlib.redirect_stderr(devnull):
                    reader = PdfReader(str(file_path))
                    _ = len(reader.pages)
            finally:
                devnull.close()
            return False, ''
        except pypdf_errors.PdfStreamError:
            return True, 'PDF 流意外结束（文件损坏或不完整）'
        except Exception as e:
            return True, f'PDF 读取失败: {type(e).__name__}'
        finally:
            logger.setLevel(old_level)
    except ImportError:
        return False, ''


def detect_broken_txt(file_path: Path) -> tuple[bool, str]:
    """检查 .txt 文件编码是否无法识别。"""
    if file_path.suffix.lower() != '.txt':
        return False, ''
    _TXT_ENCODINGS = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]
    for enc in _TXT_ENCODINGS:
        try:
            content = file_path.read_text(encoding=enc)
            if content.strip():
                return False, ''
        except (UnicodeDecodeError, UnicodeError):
            continue
    return True, f'无法用常见编码 ({"/".join(_TXT_ENCODINGS)}) 解码'


def scan_compatibility(directory: Path) -> list[tuple[Path, str]]:
    """扫描目录，返回 (问题文件路径, 原因) 列表。"""
    problems = []
    for fp in sorted(directory.rglob('*')):
        if not fp.is_file():
            continue
        ext = fp.suffix.lower()
        if ext == '.docx':
            is_prob, reason = detect_docm(fp)
        elif ext == '.pdf':
            is_prob, reason = detect_broken_pdf(fp)
        elif ext == '.txt':
            is_prob, reason = detect_broken_txt(fp)
        else:
            continue
        if is_prob:
            problems.append((fp, reason))
    return problems


def _convert_with_docling(source_path: Path, output_path: Path
                          ) -> Tuple[str, Optional[str], Optional[str]]:
    """
    用 Docling 转换文档（支持 docx/pdf）。
    返回 (status, error, warning)。
    静默 Docling 内部的图片/VML/格式不兼容等噪音。
    """
    import contextlib

    try:
        conv = _get_docling_converter()

        # 静默 Docling 内部的 stderr 噪音（图片/VML/docm 错误等）
        devnull = open(os.devnull, 'w')
        try:
            with contextlib.redirect_stderr(devnull):
                result = conv.convert(str(source_path))
        finally:
            devnull.close()

        md_content = result.document.export_to_markdown()

        if not md_content.strip():
            return ("empty", "Docling 转换内容为空", None)

        # 清洁模板段落
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "内容为空（移除了所有模板段落）", None)

        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")

        # 检查 MD 文件大小预警
        warning = _check_md_size(output_path)

        return ("ok", None, warning)

    except ImportError:
        return ("error", "缺少 docling 库", None)
    except Exception as e:
        return ("error", f"Docling {type(e).__name__}: {e}", None)


# ============================================================
# MD 文件大小预警
# ============================================================

def _check_md_size(md_path: Path) -> Optional[str]:
    """MD 文件超过阈值时返回预警信息。"""
    if md_path.exists():
        size_kb = md_path.stat().st_size / 1024
        if size_kb > MAX_MD_FILE_SIZE_KB:
            return f"MD 文件过大 ({size_kb:.0f}KB)，可能影响入库性能"
    return None


# ============================================================
# MD 内容清洁：移除封面/目录/版权/版本记录/作者信息等
# ============================================================

# 需要移除的单行模式（匹配即移除该行）
_BOILERPLATE_LINE_PATTERNS = [
    re.compile(r'^#+\s*(前\s*言|引言|概述|背景|前\s*言\s*$)'),
    re.compile(r'^#+\s*(目\s*录|目录|Contents)'),
    re.compile(r'(版权|著作权|著作权声明|版权声明|©\s*\d{4})'),
    re.compile(r'(All\s+rights?\s+reserved)', re.IGNORECASE),
    re.compile(r'(Confidential|机密|秘密|绝密)'),
    re.compile(r'(版本\s*[：:]|版\s*本\s*[：:]|版\s*本\s*号)'),
    re.compile(r'(修订记录|修\s*订\s*记\s*录|变更记录|变更历史|修订历史|文档变更)'),
    re.compile(r'^(作者|编写[：:]?|编制[：:]?|审核[：:]?|批准[：:]?|校对[：:]?|会审[：:]?|评审[：:]?|起草[：:]?|复核[：:]?)'),
    re.compile(r'^(第\s*\d+\s*页|Page\s+\d+|—\s*\d+\s*—|-\s*\d+\s*-)$'),
    re.compile(r'^中兴通讯|^ZTE\s*CORPORATION|^ZTE\s*中兴'),
    re.compile(r'^(技术文件|技术手册|产品文档|产品手册|用户手册|操作指南)'),
    re.compile(r'^文档版本\s*\d'),
    re.compile(r'^[\u2460-\u2473①-⑳]'),  # 带圈数字（常用于版权脚注）
]

# 目录检测
_TOC_ENTRY = re.compile(
    r'^\d+(\.\d+)*\s*[\u4e00-\u9fff][\u4e00-\u9fff\w]*[\s.…·]+\d+\s*$'
)
_TOC_NUM_LINE = re.compile(r'^[\d一二三四五六七八九十]+[.、．]\s*\S{1,40}$')
_TOC_PURE_DOTS = re.compile(r'^[\s.…·]+$')
_TOC_HEADING = re.compile(r'^#+\s*[目\t ]*[录録]\s*$|^[#\s]*目录|^[#\s]*Contents')


def _is_toc_section(lines: list[str], start: int, max_lookahead: int = 50) -> int:
    """检测目录段落，返回段落结束行号。"""
    if not _TOC_HEADING.match(lines[start].strip()):
        return start

    end = min(start + max_lookahead, len(lines))
    toc_count = 1
    for i in range(start + 1, end):
        raw = lines[i].strip()
        if not raw:
            toc_count += 1
            continue
        if raw.startswith('#'):
            continue
        if _TOC_PURE_DOTS.match(raw):
            toc_count += 1
            continue
        if _TOC_ENTRY.match(raw) or _TOC_NUM_LINE.match(raw):
            toc_count += 1
            continue
        break
    return start + toc_count if toc_count >= 4 else start


# 封面检测：常见封面行（短行，无编号无标题）
_COVER_LINE = re.compile(r'^[\u4e00-\u9fff\w\s]{1,30}$')
_COVER_END_MARKER = re.compile(r'^#{1,6}\s|^第[一二三四五六七八九十\d]+[章节篇]')


def _is_junk_table_row(cells: list[str]) -> bool:
    """判断表格行是否全是空单元格或纯零值（垃圾行）。"""
    for c in cells:
        s = c.strip()
        if s and s != '0':
            return False
    return True


def _compact_table_block(lines: list[str]) -> list[str]:
    """压缩表格块：去掉尾部空单元格、移除纯空/纯零行、自动计算实际列数。"""
    # 找出所有分隔线行（只检测实际表格中的分隔线 = 整行仅含 --- 和 |）
    # 并记录分隔线的列数
    sep_indices = set()
    sep_col_count = 0
    for i, line in enumerate(lines):
        cells = line.split('|')
        has_sep = sum(1 for c in cells if c.strip().replace('-', '').strip() == '' and '---' in c)
        if has_sep >= 2:
            sep_indices.add(i)
            if has_sep > sep_col_count:
                sep_col_count = has_sep

    if not sep_indices:
        # 没有分隔线：检查是否全是纯空/纯零行 → 移除
        all_junk = True
        for line in lines:
            cells = line.split('|')
            if len(cells) >= 3:
                content = [c.strip() for c in cells[1:-1]]
                if not _is_junk_table_row(content):
                    all_junk = False
                    break
        if all_junk:
            return []
        # 有内容行但仍可能有尾部空单元格，独立修剪
        out = []
        for line in lines:
            cells = line.split('|')
            if len(cells) < 3:
                out.append(line)
                continue
            cc = [c.strip() for c in cells[1:-1]]
            # 去掉尾部空单元格
            while cc and not cc[-1]:
                cc.pop()
            # 去掉尾部纯零单元格
            while cc and all(c == '0' for c in cc[-1:]):
                cc.pop()
            # 去掉前导空单元格
            while cc and not cc[0]:
                cc.pop(0)
            if _is_junk_table_row(cc):
                continue
            out.append('| ' + ' | '.join(cc) + ' |')
        return out

    if sep_col_count <= 1:
        return lines  # 单列分隔线，不处理

    # 处理每行：先去尾部空单元格，记录实际使用列数
    processed = []
    for i, line in enumerate(lines):
        cells = line.split('|')
        if len(cells) < 3:
            processed.append((i, cells, 0, False))
            continue

        content_cells = cells[1:-1]

        if i in sep_indices:
            col_count = sum(1 for c in content_cells if '---' in c)
            processed.append((i, ['---'] * col_count, col_count, True))
            continue

        # 数据行：去掉尾部空单元格
        while content_cells and not content_cells[-1].strip():
            content_cells.pop()
        while content_cells and all(c.strip() == '0' for c in content_cells[-1:]):
            content_cells.pop()

        processed.append((i, content_cells, len(content_cells), False))

    # 确定表格实际列数 = 所有数据行的最大列数（分隔线行除外）
    real_cols = max(
        (cnt for _, _, cnt, is_sep in processed if not is_sep and cnt > 0),
        default=0
    )
    if real_cols == 0:
        return []  # 纯空表（只有分隔线和空行），全部移除
    if real_cols == 1:
        return lines  # 单列表格，不处理

    out = []
    for idx, content_cells, _, is_sep in processed:
        if is_sep:
            sep_part = ['---'] * real_cols  # 强制补齐到实际列数
            out.append('|' + '|'.join(sep_part) + '|')
            continue

        if _is_junk_table_row(content_cells):
            continue

        cells = content_cells[:real_cols]
        reconstructed = '| ' + ' | '.join(c.strip() for c in cells) + ' |'
        out.append(reconstructed)

    return out


def _clean_md_content(content: str) -> str:
    """
    移除 Markdown 中的封面/目录/版权/版本记录/作者信息等模板化段落，
    以及表格中的空单元格/零值占位等冗余内容。
    策略：
      1. 先检测 TOC 段落（基于段落特征），标注范围
      2. 再移除匹配的单行模板模式
      3. 检测并移除文件开头的封面段
      4. 压缩表格：去掉尾部空单元格、移除纯空/纯零行
    """
    if not content.strip():
        return content

    lines = content.split('\n')
    n = len(lines)
    keep = [True] * n

    # Pass 1: 目录段落
    i = 0
    while i < n:
        toc_end = _is_toc_section(lines, i)
        if toc_end > i:
            for j in range(i, toc_end):
                keep[j] = False
            i = toc_end
            continue
        i += 1

    # Pass 2: 单行模板模式
    for i, line in enumerate(lines):
        if not keep[i]:
            continue
        stripped = line.strip()
        if any(p.search(stripped) for p in _BOILERPLATE_LINE_PATTERNS):
            keep[i] = False

    # Pass 3: 封面段
    first_heading_idx = -1
    for i, line in enumerate(lines):
        if not keep[i]:
            continue
        if re.match(r'^#{1,6}\s', line.strip()):
            first_heading_idx = i
            break

    if first_heading_idx > 0:
        # 第一个标题前的行 -> 如果全是短封面行 -> 移除
        pre_kept = [lines[j] for j in range(first_heading_idx) if keep[j] and lines[j].strip()]
        if pre_kept and all(len(l.strip()) < 40 for l in pre_kept):
            for j in range(first_heading_idx):
                keep[j] = False

    # Pass 4: 表格压缩（对每个表格块单独处理）
    result_lines = []
    i = 0
    while i < n:
        if not keep[i]:
            result_lines.append(lines[i])
            i += 1
            continue

        # 检测表格块：只要包含 | 且在分隔线行或空/零值行范围内的连续行
        raw = lines[i].strip()
        if '|' in raw:
            block = []
            while i < n and keep[i]:
                cur = lines[i].strip()
                if '|' not in cur:
                    break
                # 规范化：确保有前导 | 和尾随 | 方便处理
                if not cur.startswith('|'):
                    cur = '|' + cur
                if not cur.endswith('|'):
                    cur = cur + '|'
                block.append(cur)
                i += 1
            if block:
                compressed = _compact_table_block(block)
                result_lines.extend(compressed)
            continue

        result_lines.append(lines[i])
        i += 1

    cleaned = '\n'.join(result_lines)

    # ═══ 兜底核弹清洁 ═══
    #  清理 _compact_table_block 可能遗漏的表格垃圾
    #  1) 纯空行 & 纯 --- 行 → 整行删除
    cleaned = re.sub(r'^\|(?:\s*\|\s*)*\|\s*$', '', cleaned, flags=re.M)
    #  2) 空行边界合并（删除空表后可能留的多余空行）
    cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
    #  3) 每个表格行尾部空单元格截断：| a | b |   |   | → | a | b |
    #     去掉行尾连续的 | (空格) 模式
    cleaned = re.sub(r'(?<=\|)(?:\s*\|\s*)*$', '', cleaned, flags=re.M)
    #  4) 去掉前导空单元格：|   |   | a | b | → | a | b |
    cleaned = re.sub(r'^\|(?:\s*\|\s*)*\s*(?=\S)', '| ', cleaned, flags=re.M)
    #  5) 清理残留的孤立空单元格（行中连续多个空格|）  
    cleaned = re.sub(r'\|\s*\|\s*(?=\|)', '| ', cleaned)
    #  6) 移除 Docling 残留的 HTML 注释（如 <!-- image -->、<!-- 表1-1 指标集示例 --> 等）
    cleaned = re.sub(r'<!--.*?-->', '', cleaned, flags=re.DOTALL)

    cleaned = re.sub(r'\n{4,}', '\n\n\n', cleaned)
    return cleaned.strip()


# ============================================================
# 原生降级解析器
# ============================================================

def _docx_fallback(source_path: Path, output_path: Path
                   ) -> Tuple[str, Optional[str]]:
    """python-docx 降级方案（Docling 不可用时）。"""
    try:
        from docx import Document
        doc = Document(str(source_path))
        return _docx_to_md(doc, source_path, output_path)
    except ValueError as e:
        if "not a Word file" in str(e):
            return _docx_fallback_zip(source_path, output_path)
        return ("error", f"{type(e).__name__}: {e}")
    except PackageNotFoundError:
        return ("error", "不是有效的 .docx 文件（ZIP 包无法解析）")
    except KeyError as e:
        if "'NULL'" in str(e) or "NULL" in str(e):
            return _docx_fallback_zip(source_path, output_path)
        return ("error", f"{type(e).__name__}: {e}")
    except Exception as e:
        e_str = f"{type(e).__name__}: {e}"
        e_low = e_str.lower()
        if "no relationship" in e_low and "officedocument" in e_low \
           or "'null'" in e_low or "xmlsyntaxerror" in e_low \
           or ("xml" in e_low and "expected" in e_low):
            return _docx_fallback_zip(source_path, output_path)
        return ("error", e_str)


def _docx_to_md(doc, source_path: Path, output_path: Path) -> Tuple[str, Optional[str]]:
    md_lines = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style_name = para.style.name.lower() if para.style else ""
        if "heading 1" in style_name or "title" in style_name:
            md_lines.append(f"# {text}")
        elif "heading 2" in style_name:
            md_lines.append(f"## {text}")
        elif "heading 3" in style_name:
            md_lines.append(f"### {text}")
        elif "heading" in style_name:
            level = 4
            for s in ("heading 4", "heading 5", "heading 6"):
                if s in style_name:
                    level = int(s[-1])
                    break
            md_lines.append(f"{'#' * level} {text}")
        else:
            md_lines.append(text)

    for table in doc.tables:
        md_lines.append("")
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            md_lines.append("| " + " | ".join(cells) + " |")
        md_lines.append("")

    md_content = "\n\n".join(md_lines).strip()
    if not md_content:
        return ("empty", "文档内容为空")

    md_content = _clean_md_content(md_content)
    if not md_content:
        return ("empty", "文档内容为空（移除了所有模板段落）")

    _ensure_output_dir(output_path)
    output_path.write_text(md_content, encoding="utf-8")
    if not is_file_readable(output_path):
        return ("garbled", "转换后内容疑似乱码")
    return ("ok", None)


def _docx_fallback_zip(source_path: Path, output_path: Path) -> Tuple[str, Optional[str]]:
    try:
        import zipfile
        from xml.etree import ElementTree as ET
        with zipfile.ZipFile(str(source_path)) as z:
            for p in ["word/document.xml", "Word/document.xml"]:
                try:
                    xml_content = z.read(p)
                    break
                except KeyError:
                    continue
            else:
                return ("error", "ZIP 中未找到 word/document.xml")
        root = ET.fromstring(xml_content)
        ns_w = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
        md_lines = []
        for para in root.iter(f"{{{ns_w}}}p"):
            texts = [t.text for t in para.iter(f"{{{ns_w}}}t") if t.text]
            line = "".join(texts).strip()
            if line:
                md_lines.append(line)
        md_content = "\n\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "ZIP 回退提取内容为空")
        md_content = _clean_md_content(md_content)
        if not md_content:
            return ("empty", "ZIP 回退提取内容为空（移除模板后）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        return ("ok", None)
    except Exception as e:
        return ("error", f"ZIP 回退提取失败: {type(e).__name__}: {e}")


def _pdf_fallback(source_path: Path, output_path: Path
                  ) -> Tuple[str, Optional[str]]:
    """pypdf 降级方案（Docling 不可用时）。"""
    try:
        from pypdf import PdfReader, errors as pypdf_errors
        try:
            reader = PdfReader(str(source_path))
        except pypdf_errors.PdfStreamError:
            return ("error", "PDF 流意外结束（文件可能损坏或不完整）")
        except Exception as e:
            return ("error", f"PDF 读取失败: {type(e).__name__}: {e}")

        md_lines = []
        for i, page in enumerate(reader.pages, 1):
            try:
                text = page.extract_text()
            except Exception:
                text = ""
            if text and text.strip():
                md_lines.append(f"<!-- Page {i} -->\n\n{text.strip()}")

        md_content = "\n\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "PDF 内容为空（可能是扫描件）")
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "PDF 内容为空（移除了所有模板段落）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        if not is_file_readable(output_path):
            return ("garbled", "转换后内容疑似乱码")
        return ("ok", None)
    except ImportError:
        return ("error", "缺少 pypdf 库")
    except Exception as e:
        return ("error", f"PyPDF {type(e).__name__}: {e}")


def _xlsx_fallback(source_path: Path, output_path: Path
                   ) -> Tuple[str, Optional[str]]:
    """openpyxl 降级方案（Docling 不可用时）。"""
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(source_path), read_only=True, data_only=True)
        md_lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            md_lines.append(f"## {sheet_name}")
            md_lines.append("")
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
            headers = [str(c) if c is not None else "" for c in rows[0]]
            md_lines.append("| " + " | ".join(headers) + " |")
            md_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            data_rows = [row for row in rows[1:] if not _is_junk_xlsx_row(row)]
            for row in data_rows:
                cells = [str(c) if c is not None else "" for c in row]
                md_lines.append("| " + " | ".join(cells) + " |")
            md_lines.append("")
        wb.close()
        md_content = "\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "Excel 内容为空")
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "Excel 内容为空（移除了所有模板段落）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        return ("ok", None)
    except ImportError:
        return ("error", "缺少 openpyxl 库")
    except Exception as e:
        return ("error", f"{type(e).__name__}: {e}")


def _pptx_fallback(source_path: Path, output_path: Path
                   ) -> Tuple[str, Optional[str]]:
    """python-pptx 降级方案（Docling 不可用时）。"""
    try:
        from pptx import Presentation
        prs = Presentation(str(source_path))
        md_lines = []
        for i, slide in enumerate(prs.slides, 1):
            md_lines.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            md_lines.append(text)
            md_lines.append("")
        md_content = "\n".join(md_lines).strip()
        if not md_content:
            return ("empty", "幻灯片内容为空")
        md_content = _clean_md_content(md_content)
        if not md_content.strip():
            return ("empty", "幻灯片内容为空（移除了所有模板段落）")
        _ensure_output_dir(output_path)
        output_path.write_text(md_content, encoding="utf-8")
        return ("ok", None)
    except ImportError:
        return ("error", "缺少 python-pptx 库")
    except Exception as e:
        return ("error", f"{type(e).__name__}: {e}")


# ============================================================
# xlsx 空行/零值行过滤（供降级用）
# ============================================================

def _is_junk_cell(value) -> bool:
    if value is None:
        return True
    if isinstance(value, (int, float)):
        return value == 0 or value == 0.0
    s = str(value).strip()
    return s == "" or s == "0" or s == "0.0"


def _is_junk_xlsx_row(row: tuple) -> bool:
    if not row:
        return True
    return all(_is_junk_cell(cell) for cell in row)


# ============================================================
# Markdown (.md) — 直接复制
# ============================================================

def _convert_md(source_path: Path, output_path: Path) -> Tuple[str, Optional[str], Optional[str]]:
    try:
        content = source_path.read_text(encoding="utf-8")
        if not content.strip():
            return ("empty", "文件内容为空", None)
        if not is_file_readable(source_path):
            return ("garbled", "源文件内容疑似乱码", None)
        content = _clean_md_content(content)
        if not content.strip():
            return ("empty", "文件内容为空（移除了所有模板段落）", None)
        _ensure_output_dir(output_path)
        output_path.write_text(content, encoding="utf-8")
        warning = _check_md_size(output_path)
        return ("ok", None, warning)
    except Exception as e:
        return ("error", f"{type(e).__name__}: {e}", None)


# ============================================================
# 纯文本 (.txt) — 支持编码回退
# ============================================================

_TXT_ENCODINGS = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]


def _convert_txt(source_path: Path, output_path: Path) -> Tuple[str, Optional[str], Optional[str]]:
    for enc in _TXT_ENCODINGS:
        try:
            content = source_path.read_text(encoding=enc)
            if content.strip():
                break
        except (UnicodeDecodeError, UnicodeError):
            continue
    else:
        return ("error", f"无法用 {', '.join(_TXT_ENCODINGS)} 解码文件内容", None)
    if not content.strip():
        return ("empty", "文件内容为空", None)
    content = _clean_md_content(content)
    if not content.strip():
        return ("empty", "文件内容为空（移除了所有模板段落）", None)
    _ensure_output_dir(output_path)
    output_path.write_text(content, encoding="utf-8")
    warning = _check_md_size(output_path)
    return ("ok", None, warning)


# ============================================================
# 原生降级调度器
# ============================================================

_FALLBACK_REGISTRY = {
    ".docx": _docx_fallback,
    ".pdf":  _pdf_fallback,
    ".xlsx": _xlsx_fallback,
    ".pptx": _pptx_fallback,
}


def convert_single_file(source_path: Path) -> dict:
    """
    转换单个源文件为 Markdown。

    Parameters
    ----------
    source_path : Path
        源文件绝对路径。

    Returns
    -------
    dict: {
        "rel_path": ...,
        "status": "ok" | "garbled" | "empty" | "error" | "skip",
        "md_path": ... | None,
        "error": ... | None,
        "warning": ... | None,
        "sha256": ...,
        "size": ...,
        "mtime": ...,
    }
    """
    rel_path = _get_rel_path(source_path)
    ext = source_path.suffix.lower()
    result = {
        "rel_path": rel_path,
        "status": "error",
        "md_path": None,
        "error": None,
        "warning": None,
        "sha256": "",
        "size": 0,
        "mtime": "",
    }

    try:
        stat = source_path.stat()
        result["size"] = stat.st_size
        result["sha256"] = compute_sha256(source_path)

        from datetime import datetime, timezone
        dt = datetime.fromtimestamp(stat.st_mtime, tz=timezone.utc)
        result["mtime"] = dt.isoformat(timespec="seconds")

        if ext not in SUPPORTED_EXTENSIONS:
            result["status"] = "skip"
            result["error"] = f"不支持格式 {ext}，忽略: {rel_path}"
            return result

        # 检测伪装为 .docx 的宏文档 → 跳过，需手工转换
        if ext == ".docx" and _is_docm_file(source_path):
            result["status"] = "skip"
            result["error"] = f"宏文档 (.docm)，需用 Word 另存为 .docx: {rel_path}"
            return result

        output_path = _get_output_md_path(source_path)

        # ── 策略：Docling 主力(docx/pdf) → 原生降级(xlsx/pptx) → 直接复制(md/txt) ──
        if ext in _DOCLING_SUPPORTED:
            # 主力：Docling
            status, error, warning = _convert_with_docling(source_path, output_path)
            if status == "ok":
                result["status"] = status
                result["error"] = error
                result["warning"] = warning
            else:
                # Docling 失败 → 降级到原生解析器
                fallback = _FALLBACK_REGISTRY.get(ext)
                if fallback:
                    status_fb, error_fb = fallback(source_path, output_path)
                    if status_fb == "ok":
                        result["status"] = status_fb
                        result["error"] = None
                        result["warning"] = _check_md_size(output_path)
                    else:
                        result["status"] = status_fb
                        result["error"] = error_fb
                else:
                    result["status"] = status
                    result["error"] = error
        elif ext in _FALLBACK_REGISTRY:
            # xlsx/pptx：直接使用原生解析器
            fallback = _FALLBACK_REGISTRY[ext]
            status_fb, error_fb = fallback(source_path, output_path)
            result["status"] = status_fb
            result["error"] = error_fb
            if status_fb == "ok":
                result["warning"] = _check_md_size(output_path)
        elif ext == ".md":
            status, error, warning = _convert_md(source_path, output_path)
            result["status"] = status
            result["error"] = error
            result["warning"] = warning
        elif ext == ".txt":
            status, error, warning = _convert_txt(source_path, output_path)
            result["status"] = status
            result["error"] = error
            result["warning"] = warning

        if result["status"] == "ok":
            result["md_path"] = str(output_path.with_suffix(".md").absolute())

    except Exception as e:
        result["status"] = "error"
        result["error"] = f"意外异常: {type(e).__name__}: {e}\n{traceback.format_exc()}"

    return result


def convert_batch(file_paths: List[Path],
                  max_workers: int = CONVERT_WORKERS,
                  progress_callback=None) -> List[dict]:
    """批量转换文件，支持 Ctrl+C 中断。"""
    results = []
    with ThreadPoolExecutor(max_workers=max_workers) as pool:
        futures = {pool.submit(convert_single_file, fp): fp for fp in file_paths}
        try:
            for future in as_completed(futures):
                result = future.result()
                results.append(result)
                if progress_callback:
                    progress_callback(result)
        except KeyboardInterrupt:
            for f in futures:
                f.cancel()
            pool.shutdown(wait=False)
            raise

    results.sort(key=lambda r: r["rel_path"])
    return results
